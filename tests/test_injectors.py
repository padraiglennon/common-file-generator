"""End-to-end injector tests: inject, re-open the output, assert content."""

from __future__ import annotations

from pathlib import Path

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from common_file_generator.core import Config, MediaSpec, TableSpec, generate


def _config(image: Path) -> Config:
    return Config(
        text={"title": "Q3 Review"},
        tables=[
            TableSpec(
                name="RevenueTable",
                header=["Region", "Q2", "Q3"],
                rows=[["North", "120", "138"], ["South", "98", "101"]],
            )
        ],
        media=[MediaSpec(name="Logo", path=image)],
    )


def test_pptx_injection(
    pptx_template: Path, sample_image: Path, tmp_path: Path
) -> None:
    out = tmp_path / "out.pptx"
    report = generate(pptx_template, _config(sample_image), out)
    assert not report.has_problems
    assert out.is_file()

    prs = Presentation(str(out))
    shapes = list(prs.slides[0].shapes)
    texts = [s.text_frame.text for s in shapes if s.has_text_frame]
    assert "Q3 Review" in texts

    table = next(s.table for s in shapes if s.has_table)
    assert table.cell(0, 0).text == "Region"
    assert table.cell(1, 0).text == "North"
    assert any(s.shape_type == 13 for s in shapes)  # 13 == PICTURE


def test_docx_injection(
    docx_template: Path, sample_image: Path, tmp_path: Path
) -> None:
    out = tmp_path / "out.docx"
    report = generate(docx_template, _config(sample_image), out)
    assert not report.has_problems

    doc = Document(str(out))
    assert any(p.text == "Q3 Review" for p in doc.paragraphs)
    table = doc.tables[0]
    assert table.cell(0, 0).text == "Region"
    assert table.cell(1, 0).text == "North"
    assert doc.inline_shapes  # the logo was inserted


def test_xlsx_injection(
    xlsx_template: Path, sample_image: Path, tmp_path: Path
) -> None:
    out = tmp_path / "out.xlsx"
    report = generate(xlsx_template, _config(sample_image), out)
    assert not report.has_problems

    wb = load_workbook(str(out))
    ws = wb["Data"]
    assert ws["A1"].value == "Q3 Review"
    assert ws["A3"].value == "Region"
    assert ws["A4"].value == "North"
    assert ws._images  # the logo was anchored


def test_report_written_next_to_output(
    pptx_template: Path, sample_image: Path, tmp_path: Path
) -> None:
    out = tmp_path / "out.pptx"
    generate(pptx_template, _config(sample_image), out)
    assert (tmp_path / "out.report.txt").is_file()
