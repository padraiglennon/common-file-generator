"""Tests for complexity-driven generate mode (designed slide types)."""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.enum.dml import MSO_FILL

from ms_office_file_generator.core import Complexity, generate_deck
from ms_office_file_generator.generators import PptxComplexityGenerator

_SLIDE_W = 12192000  # 13.333 in, in EMU
_SLIDE_H = 6858000  # 7.5 in, in EMU
_TOL = 12700  # 1pt tolerance


def test_slide_count_matches_request(tmp_path: Path) -> None:
    out = tmp_path / "deck.pptx"
    generate_deck(str(out), complexity="standard", slides=12, seed=1)
    assert len(Presentation(str(out)).slides) == 12


def test_invalid_slide_count_raises() -> None:
    with pytest.raises(ValueError, match="at least 1"):
        PptxComplexityGenerator(Complexity.MINIMAL, slides=0)


def test_no_shape_extends_past_the_slide(tmp_path: Path) -> None:
    for level in ("minimal", "simple", "standard", "complex", "maximum"):
        out = tmp_path / f"{level}.pptx"
        generate_deck(str(out), complexity=level, slides=15, seed=4)
        prs = Presentation(str(out))
        for n, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if None in (shape.left, shape.top, shape.width, shape.height):
                    continue
                assert shape.left + shape.width <= _SLIDE_W + _TOL, (
                    f"{level} slide {n}: shape runs off the right edge"
                )
                assert shape.top + shape.height <= _SLIDE_H + _TOL, (
                    f"{level} slide {n}: shape runs off the bottom edge"
                )


def test_tables_never_overflow_the_slide_bottom(tmp_path: Path) -> None:
    out = tmp_path / "max.pptx"
    generate_deck(str(out), complexity="maximum", slides=40, seed=2)
    prs = Presentation(str(out))
    tables_seen = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            tables_seen += 1
            rendered = sum(row.height for row in shape.table.rows)
            assert shape.top + rendered <= _SLIDE_H + _TOL
    assert tables_seen > 0  # the run actually exercised tables


def test_minimal_uses_only_bullets(tmp_path: Path) -> None:
    out = tmp_path / "min.pptx"
    generate_deck(str(out), complexity="minimal", slides=10, seed=5)
    prs = Presentation(str(out))
    content = [s for i, slide in enumerate(prs.slides) if i > 0 for s in slide.shapes]
    assert not any(s.has_chart for s in content)
    assert not any(s.has_table for s in content)


def test_maximum_video_slide_links_to_youtube(tmp_path: Path) -> None:
    out = tmp_path / "max.pptx"
    url = "https://www.youtube.com/watch?v=abc123"
    generate_deck(str(out), complexity="maximum", slides=40, seed=2, video_url=url)
    prs = Presentation(str(out))
    found = _collect_hyperlinks(prs)
    assert url in found


def test_image_bullets_slides_have_image_and_bullets(tmp_path: Path) -> None:
    out = tmp_path / "max.pptx"
    generate_deck(str(out), complexity="maximum", slides=60, seed=0)
    prs = Presentation(str(out))
    matches = 0
    for slide in prs.slides:
        has_pic = any(s.shape_type == 13 for s in slide.shapes)
        has_bullets = any(
            s.has_text_frame and "▸" in s.text_frame.text for s in slide.shapes
        )
        if has_pic and has_bullets:
            matches += 1
    assert matches > 0


def test_theme_background_is_consistent(tmp_path: Path) -> None:
    out = tmp_path / "bg.pptx"
    generate_deck(str(out), complexity="simple", slides=8, seed=1, background="theme")
    prs = Presentation(str(out))
    # Every slide, including the title slide, gets the consistent tint.
    tints = {str(slide.background.fill.fore_color.rgb) for slide in prs.slides}
    assert tints == {"EAF2F8"}


def test_random_background_varies(tmp_path: Path) -> None:
    out = tmp_path / "bg.pptx"
    generate_deck(str(out), complexity="simple", slides=12, seed=1, background="random")
    prs = Presentation(str(out))
    tints = {
        str(slide.background.fill.fore_color.rgb)
        for i, slide in enumerate(prs.slides)
        if i > 0
    }
    assert len(tints) > 1


def test_explicit_background_color(tmp_path: Path) -> None:
    out = tmp_path / "bg.pptx"
    generate_deck(
        str(out), complexity="simple", slides=5, seed=1, background_color="FFEFD5"
    )
    prs = Presentation(str(out))
    # Includes the title slide (slide 0).
    tints = {str(slide.background.fill.fore_color.rgb) for slide in prs.slides}
    assert tints == {"FFEFD5"}


def test_invalid_background_color_raises(tmp_path: Path) -> None:
    with pytest.raises(ValueError, match="Invalid background colour"):
        generate_deck(str(tmp_path / "x.pptx"), slides=3, background_color="nothex")


_DIVIDER_BG = "2B3A42"  # OCEAN.divider_background


def _is_divider(slide: object) -> bool:
    fill = slide.background.fill
    if fill.type != MSO_FILL.SOLID:
        return False
    return str(fill.fore_color.rgb) == _DIVIDER_BG


def test_dividers_stay_rare(tmp_path: Path) -> None:
    # Sparse divider slides must be punctuation, not a large share of the deck.
    out = tmp_path / "max.pptx"
    generate_deck(str(out), complexity="maximum", slides=70, seed=0)
    prs = Presentation(str(out))
    dividers = sum(
        1 for i, slide in enumerate(prs.slides) if i > 0 and _is_divider(slide)
    )
    assert dividers <= len(list(prs.slides)) * 0.15


def test_no_content_slide_is_title_only(tmp_path: Path) -> None:
    # Every content slide must carry body content, not just a heading - a
    # title-only slide reads as empty. Dividers included.
    for level in ("minimal", "simple", "standard", "complex", "maximum"):
        out = tmp_path / f"{level}.pptx"
        generate_deck(str(out), complexity=level, slides=60, seed=0)
        prs = Presentation(str(out))
        for n, slide in enumerate(prs.slides):
            if n == 0:
                continue  # the title slide is allowed to be a title + subtitle
            assert len(list(slide.shapes)) >= 2, (
                f"{level} slide {n} is title-only (no content)"
            )


def test_divider_slide_has_subtitle(tmp_path: Path) -> None:
    out = tmp_path / "max.pptx"
    generate_deck(str(out), complexity="maximum", slides=70, seed=0)
    prs = Presentation(str(out))
    # A divider carries a heading, an accent shape, and a supporting subtitle.
    dividers = [
        slide for i, slide in enumerate(prs.slides) if i > 0 and _is_divider(slide)
    ]
    assert dividers, "no divider slide was drawn to exercise this test"
    for slide in dividers:
        assert any(s.shape_type == 1 for s in slide.shapes)  # AUTO_SHAPE accent
        texts = [
            s.text_frame.text
            for s in slide.shapes
            if s.shape_type == 17 and s.text_frame.text.strip()  # TEXT_BOX
        ]
        assert len(texts) >= 2  # heading + subtitle, both non-empty


def test_seed_is_deterministic(tmp_path: Path) -> None:
    a, b = tmp_path / "a.pptx", tmp_path / "b.pptx"
    generate_deck(str(a), complexity="maximum", slides=12, seed=42)
    generate_deck(str(b), complexity="maximum", slides=12, seed=42)
    assert _slide_texts(a) == _slide_texts(b)


def test_different_seed_differs(tmp_path: Path) -> None:
    a, b = tmp_path / "a.pptx", tmp_path / "b.pptx"
    generate_deck(str(a), complexity="maximum", slides=12, seed=1)
    generate_deck(str(b), complexity="maximum", slides=12, seed=2)
    assert _slide_texts(a) != _slide_texts(b)


def test_missing_theme_raises(tmp_path: Path) -> None:
    with pytest.raises(FileNotFoundError, match="Theme template not found"):
        generate_deck(
            str(tmp_path / "out.pptx"),
            complexity="standard",
            slides=3,
            theme_path=str(tmp_path / "nope.pptx"),
        )


def _collect_hyperlinks(prs: Presentation) -> set[str]:
    found: set[str] = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            action = getattr(shape, "click_action", None)
            if action is not None and action.hyperlink.address:
                found.add(action.hyperlink.address)
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.hyperlink.address:
                            found.add(run.hyperlink.address)
    return found


def _slide_texts(path: Path) -> list[str]:
    prs = Presentation(str(path))
    return [
        shape.text_frame.text
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_text_frame
    ]
