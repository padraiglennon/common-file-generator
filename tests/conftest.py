"""Shared fixtures: build minimal Golden templates and a sample image on the fly.

Templates are generated in code (not committed binaries) so the tests are
self-contained and reviewable. The generation here is *test scaffolding*, not the
product - the product never creates layout.
"""

from __future__ import annotations

from pathlib import Path

import pytest
from docx import Document
from openpyxl import Workbook
from openpyxl.workbook.defined_name import DefinedName
from PIL import Image
from pptx import Presentation
from pptx.util import Inches


@pytest.fixture
def sample_image(tmp_path: Path) -> Path:
    path = tmp_path / "logo.png"
    Image.new("RGB", (16, 16), "red").save(path)
    return path


@pytest.fixture
def pptx_template(tmp_path: Path) -> Path:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(6), Inches(1))
    title.text_frame.paragraphs[0].add_run().text = "{{title}}"

    table_shape = slide.shapes.add_table(
        3, 3, Inches(1), Inches(2), Inches(6), Inches(2)
    )
    table_shape.name = "RevenueTable"

    logo = slide.shapes.add_textbox(Inches(7), Inches(0.5), Inches(1), Inches(1))
    logo.name = "Logo"

    path = tmp_path / "template.pptx"
    prs.save(path)
    return path


@pytest.fixture
def docx_template(tmp_path: Path) -> Path:
    doc = Document()
    doc.add_paragraph("{{title}}")
    table = doc.add_table(rows=3, cols=3)
    table.cell(2, 2).text = "{{table:RevenueTable}}"
    doc.add_paragraph("{{media:Logo}}")
    path = tmp_path / "template.docx"
    doc.save(path)
    return path


@pytest.fixture
def xlsx_template(tmp_path: Path) -> Path:
    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    # Reserve a grid so anchors have room to fill into.
    for row in range(1, 11):
        for col in range(1, 6):
            ws.cell(row=row, column=col, value="")
    wb.defined_names.add(DefinedName("title", attr_text="Data!$A$1"))
    wb.defined_names.add(DefinedName("RevenueTable", attr_text="Data!$A$3"))
    wb.defined_names.add(DefinedName("Logo", attr_text="Data!$E$1"))
    path = tmp_path / "template.xlsx"
    wb.save(path)
    return path
