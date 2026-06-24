"""Complexity-driven PowerPoint generator.

Builds a deck of N slides by randomly drawing **designed slide types** from a
pool determined by the complexity level. Each slide has a single focal element
(bullets, a chevron list, an image+text pane, a chart, a fit-to-slot table, a
video link, or a section divider), so slides never overflow or overlap and the
deck mixes layouts the way a real one does.

Content is seeded lorem-ipsum, so the same (slides, complexity, seed) always
yields the same deck. Placeholder images and the video poster are generated in
memory, so a deck needs no external assets.

When ``theme_path`` is given, that designed ``.pptx`` is used as the base so its
slide-master backgrounds, colours, and layouts carry through; otherwise the
in-code :data:`~.theme.DEFAULT_THEME` is applied.
"""

from __future__ import annotations

import random
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from ms_office_file_generator.core.complexity import Complexity, slide_pool
from ms_office_file_generator.core.lorem import Lorem
from ms_office_file_generator.generators import slides as slide_types
from ms_office_file_generator.generators import style
from ms_office_file_generator.generators.background import Background
from ms_office_file_generator.generators.slides import SlideContext
from ms_office_file_generator.generators.theme import DEFAULT_THEME

_SLIDE_W = Inches(13.333)
_SLIDE_H = Inches(7.5)
_BLANK_LAYOUT = 6

# Maps a pool key to (builder, takes_index).
_BUILDERS = {
    "bullets": slide_types.bullets_slide,
    "chevron": slide_types.chevron_list_slide,
    "image_text": slide_types.image_text_slide,
    "image_bullets": slide_types.image_bullets_slide,
    "chart": slide_types.chart_slide,
    "table": slide_types.table_slide,
    "video": slide_types.video_slide,
    "divider": slide_types.divider_slide,
}


class PptxComplexityGenerator:
    """Generate a designed PowerPoint deck at a chosen complexity level."""

    def __init__(
        self,
        complexity: Complexity = Complexity.STANDARD,
        *,
        slides: int = 10,
        seed: int = 0,
        video_url: str = "",
        theme_path: str | None = None,
        background: Background | None = None,
    ) -> None:
        if slides < 1:
            raise ValueError("slides must be at least 1")
        self.complexity = complexity
        self.slides = slides
        self.theme_path = theme_path
        self._background = background or Background(DEFAULT_THEME)
        pool = slide_pool(complexity)
        self._pool_keys = [key for key, _ in pool]
        self._pool_weights = [weight for _, weight in pool]
        self._rng = random.Random(seed)
        self._ctx = SlideContext(
            rng=self._rng,
            lorem=Lorem(self._rng),
            theme=DEFAULT_THEME,
            video_url=video_url,
        )

    def build(self) -> Presentation:
        prs = self._base_presentation()
        self._add_title_slide(prs)
        for index in range(1, self.slides):
            self._add_content_slide(prs, index)
        return prs

    def save(self, out_path: str | Path) -> Path:
        out_path = Path(out_path)
        out_path.parent.mkdir(parents=True, exist_ok=True)
        self.build().save(str(out_path))
        return out_path

    # --- internals ------------------------------------------------------

    def _base_presentation(self) -> Presentation:
        if self.theme_path is not None:
            path = Path(self.theme_path)
            if not path.is_file():
                raise FileNotFoundError(f"Theme template not found: {path}")
            return Presentation(str(path))
        prs = Presentation()
        prs.slide_width = _SLIDE_W
        prs.slide_height = _SLIDE_H
        return prs

    def _add_title_slide(self, prs: Presentation) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide_types.title_slide(slide, self._ctx)

    def _add_content_slide(self, prs: Presentation, index: int) -> None:
        layout = prs.slide_layouts[_BLANK_LAYOUT]
        slide = prs.slides.add_slide(layout)
        color = self._background.color_for(self._rng)
        if color is not None:
            style.set_background(slide, color)
        key = self._rng.choices(self._pool_keys, weights=self._pool_weights)[0]
        _BUILDERS[key](slide, self._ctx, index)
