"""PowerPoint injector (python-pptx).

Populates ``{{token}}`` text, named tables, and named picture placeholders in a
Golden ``.pptx`` template. No slides or shapes are created; only existing
elements are filled.
"""

from __future__ import annotations

from collections.abc import Iterator
from pathlib import Path

from pptx import Presentation
from pptx.shapes.base import BaseShape

from common_file_generator.core.config import Config, MediaSpec, TableSpec
from common_file_generator.core.injector import Injector
from common_file_generator.injectors._text import (
    find_tokens,
    replace_in_paragraph,
)


class PptxInjector(Injector):
    """Inject data into a PowerPoint Golden template."""

    extensions = frozenset({".pptx"})

    def _load(self, path: Path) -> Presentation:
        return Presentation(str(path))

    def _save(self, path: Path) -> None:
        self._document.save(str(path))

    def _iter_shapes(self) -> Iterator[BaseShape]:
        for slide in self._document.slides:
            yield from slide.shapes

    # --- text -----------------------------------------------------------

    def _inject_text(self, config: Config) -> None:
        if not config.text:
            return
        seen: set[str] = set()
        for shape in self._iter_shapes():
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                seen |= find_tokens(paragraph.text)
                replace_in_paragraph(paragraph, config.text)

        for name in sorted(set(config.text) - seen):
            self.report.warning(
                f"text.{name}",
                f"no '{{{{{name}}}}}' placeholder was found in the template; "
                "this value was not used.",
            )

    # --- tables ---------------------------------------------------------

    def _inject_tables(self, config: Config) -> None:
        if not config.tables:
            return
        tables_by_name = {
            shape.name: shape for shape in self._iter_shapes() if shape.has_table
        }
        for spec in config.tables:
            shape = tables_by_name.get(spec.name)
            if shape is None:
                self.report.error(
                    f"tables.{spec.name}",
                    "no table with this name was found in the template.",
                )
                continue
            self._fill_table(shape.table, spec)

    def _fill_table(self, table: object, spec: TableSpec) -> None:
        rows = table.rows  # type: ignore[attr-defined]
        cols = table.columns  # type: ignore[attr-defined]
        n_rows, n_cols = len(rows), len(cols)

        data_rows = list(spec.rows)
        start = 0
        if spec.header is not None:
            data_rows = [spec.header, *data_rows]

        if len(data_rows) > n_rows:
            self.report.warning(
                f"tables.{spec.name}",
                f"you provided {len(data_rows)} rows but the template table holds "
                f"{n_rows}; kept the first {n_rows}.",
            )
        for r, row in enumerate(data_rows[:n_rows], start=start):
            if len(row) > n_cols:
                self.report.warning(
                    f"tables.{spec.name}",
                    f"row {r} has {len(row)} cells but the table has {n_cols} "
                    f"columns; kept the first {n_cols}.",
                )
            for c, value in enumerate(row[:n_cols]):
                table.cell(r, c).text = value  # type: ignore[attr-defined]

    # --- media ----------------------------------------------------------

    def _inject_media(self, config: Config) -> None:
        if not config.media:
            return
        shapes_by_name = {shape.name: shape for shape in self._iter_shapes()}
        for spec in config.media:
            self._place_media(shapes_by_name, spec)

    def _place_media(
        self, shapes_by_name: dict[str, BaseShape], spec: MediaSpec
    ) -> None:
        shape = shapes_by_name.get(spec.name)
        if shape is None:
            self.report.error(
                f"media.{spec.name}",
                "no shape with this name was found in the template.",
            )
            return
        if not spec.path.is_file():
            self.report.error(
                f"media.{spec.name}",
                f"image file not found: {spec.path}.",
            )
            return

        slide_shapes = shape._parent  # the SlideShapes collection
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
        shape._element.getparent().remove(shape._element)
        slide_shapes.add_picture(str(spec.path), left, top, width, height)
