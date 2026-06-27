"""Low-level styling helpers shared by the designed slide types.

These wrap the python-pptx fill/font/line calls so the slide-type code stays
readable. Everything is driven by a :class:`~.theme.Theme`.
"""

from __future__ import annotations

import contextlib

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

from common_file_generator.generators.layout import Box
from common_file_generator.generators.theme import Theme


def set_background(slide: object, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_heading(slide: object, box: Box, text: str, theme: Theme) -> object:
    shape = slide.shapes.add_textbox(box.left, box.top, box.width, box.height)
    frame = shape.text_frame
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.name = theme.title_font
    run.font.color.rgb = theme.heading_color
    return shape


def add_body(
    slide: object, box: Box, text: str, theme: Theme, *, size: int = 16
) -> None:
    shape = slide.shapes.add_textbox(box.left, box.top, box.width, box.height)
    frame = shape.text_frame
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.name = theme.body_font
    run.font.color.rgb = theme.body_color


def add_gradient_bar(
    slide: object,
    box: Box,
    text: str,
    theme: Theme,
    *,
    ratio: float = 0.5,
) -> None:
    """A rounded rectangle with a top-to-bottom gradient and centred white text.

    ``ratio`` (0..1) shifts the gradient between the theme's start and end colour
    so a stacked list fades from one to the other.
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, box.left, box.top, box.width, box.height
    )
    fill = shape.fill
    fill.gradient()
    blended = _blend(theme.gradient_start, theme.gradient_end, ratio)
    blended_2 = _blend(theme.gradient_start, theme.gradient_end, min(1.0, ratio + 0.15))
    fill.gradient_stops[0].color.rgb = blended
    fill.gradient_stops[1].color.rgb = blended_2
    with contextlib.suppress(NotImplementedError, ValueError):
        fill.gradient_angle = 90.0
    shape.line.fill.background()

    frame = shape.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.name = theme.body_font
    run.font.color.rgb = theme.on_fill_color


def style_table(
    table: object, theme: Theme, *, font_size: int, row_height: Emu
) -> None:
    """Apply themed header colour, small font, and tight margins to a table."""
    for r, row in enumerate(table.rows):
        row.height = row_height
        for cell in row.cells:
            cell.margin_top = Pt(1)
            cell.margin_bottom = Pt(1)
            cell.margin_left = Pt(4)
            cell.margin_right = Pt(4)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(font_size)
                    run.font.name = theme.body_font
                    if r == 0:
                        run.font.bold = True
                        run.font.color.rgb = theme.on_fill_color


def _blend(a: RGBColor, b: RGBColor, ratio: float) -> RGBColor:
    ratio = max(0.0, min(1.0, ratio))
    return RGBColor(
        round(a[0] + (b[0] - a[0]) * ratio),
        round(a[1] + (b[1] - a[1]) * ratio),
        round(a[2] + (b[2] - a[2]) * ratio),
    )
