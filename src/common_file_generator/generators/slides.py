"""Designed slide types, each with a single focal element.

Every builder takes a :class:`SlideContext` (RNG, lorem, theme, helpers) and a
python-pptx ``slide`` and lays out one designed slide. Keeping one focal element
per slide is what avoids the overflow/overlap problems of cramming many large
elements into a grid, and lets the deck mix layouts the way a real deck does.
"""

from __future__ import annotations

import io
import random
from dataclasses import dataclass

from PIL import Image, ImageDraw
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from common_file_generator.core.lorem import Lorem
from common_file_generator.generators import style
from common_file_generator.generators.layout import Box
from common_file_generator.generators.theme import Theme

_CHART_TYPES = (
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    XL_CHART_TYPE.LINE_MARKERS,
    XL_CHART_TYPE.PIE,
)
_HEADING = Box(Inches(0.6), Inches(0.4), Inches(12.1), Inches(1.0))
_FULL_BODY = Box(Inches(0.6), Inches(1.6), Inches(12.1), Inches(5.4))


@dataclass
class SlideContext:
    """Everything a slide builder needs, passed in instead of a god-object."""

    rng: random.Random
    lorem: Lorem
    theme: Theme
    video_url: str

    def heading_text(self, index: int) -> str:
        return f"{index}. {self.lorem.title()}"

    def placeholder_image(self, w: int = 480, h: int = 300) -> io.BytesIO:
        color = (
            self.rng.randint(40, 220),
            self.rng.randint(40, 220),
            self.rng.randint(40, 220),
        )
        return _png(Image.new("RGB", (w, h), color))


# --- slide builders -----------------------------------------------------


def title_slide(slide: object, ctx: SlideContext) -> None:
    title = slide.shapes.title
    if title is not None:
        title.text = ctx.lorem.title()
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = ctx.lorem.sentence()


def divider_slide(slide: object, ctx: SlideContext, index: int) -> None:
    style.set_background(slide, ctx.theme.divider_background)
    box = Box(Inches(1.2), Inches(2.4), Inches(11), Inches(1.6))
    shape = slide.shapes.add_textbox(box.left, box.top, box.width, box.height)
    frame = shape.text_frame
    frame.word_wrap = True
    run = frame.paragraphs[0].add_run()
    run.text = ctx.lorem.title()
    run.font.size = Pt(40)
    run.font.name = ctx.theme.title_font
    run.font.color.rgb = ctx.theme.accent

    accent = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.25),
        Inches(4.0),
        Inches(2.2),
        Inches(0.12),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = ctx.theme.accent
    accent.line.fill.background()

    subtitle = slide.shapes.add_textbox(
        Inches(1.2), Inches(4.3), Inches(11), Inches(1.4)
    )
    sub_frame = subtitle.text_frame
    sub_frame.word_wrap = True
    sub_run = sub_frame.paragraphs[0].add_run()
    sub_run.text = ctx.lorem.sentence(8, 16)
    sub_run.font.size = Pt(18)
    sub_run.font.name = ctx.theme.body_font
    sub_run.font.color.rgb = ctx.theme.on_fill_color


def bullets_slide(slide: object, ctx: SlideContext, index: int) -> None:
    style.add_heading(slide, _HEADING, ctx.heading_text(index), ctx.theme)
    _add_bullets(slide, ctx, _FULL_BODY, count=ctx.rng.randint(4, 7), size=18)


def _add_bullets(
    slide: object, ctx: SlideContext, box: Box, *, count: int, size: int
) -> None:
    shape = slide.shapes.add_textbox(box.left, box.top, box.width, box.height)
    frame = shape.text_frame
    frame.word_wrap = True
    for i, text in enumerate(ctx.lorem.bullets(count)):
        paragraph = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        paragraph.text = f"▸  {text}"
        paragraph.font.size = Pt(size)
        paragraph.font.name = ctx.theme.body_font
        paragraph.font.color.rgb = ctx.theme.body_color
        paragraph.space_after = Pt(12)


def chevron_list_slide(slide: object, ctx: SlideContext, index: int) -> None:
    style.add_heading(slide, _HEADING, ctx.heading_text(index), ctx.theme)
    count = ctx.rng.randint(4, 7)
    top0, bar_h, gap = 1.7, (5.2 / count) - 0.12, 0.12
    width, left = 11.0, 1.2
    for i in range(count):
        top = top0 + i * (bar_h + gap)
        box = Box(Inches(left), Inches(top), Inches(width), Inches(bar_h))
        ratio = i / max(1, count - 1)
        style.add_gradient_bar(
            slide, box, ctx.lorem.sentence(3, 8), ctx.theme, ratio=ratio
        )


def image_text_slide(slide: object, ctx: SlideContext, index: int) -> None:
    # Left text pane, right image - the two-pane look.
    style.add_heading(
        slide,
        Box(Inches(0.6), Inches(0.5), Inches(5.6), Inches(1.4)),
        ctx.lorem.title(4),
        ctx.theme,
    )
    style.add_body(
        slide,
        Box(Inches(0.6), Inches(2.0), Inches(5.6), Inches(4.8)),
        ctx.lorem.paragraph(3, 5),
        ctx.theme,
        size=18,
    )
    image = ctx.placeholder_image(720, 900)
    slide.shapes.add_picture(image, Inches(6.7), Inches(0.0), Inches(6.63), Inches(7.5))


def image_bullets_slide(slide: object, ctx: SlideContext, index: int) -> None:
    """Image pane on one (random) side, a bullet list on the other."""
    style.add_heading(slide, _HEADING, ctx.heading_text(index), ctx.theme)
    image_left = ctx.rng.random() < 0.5
    image_pane = Box(
        Inches(0.6 if image_left else 6.9), Inches(1.7), Inches(5.8), Inches(5.1)
    )
    text_pane = Box(
        Inches(6.9 if image_left else 0.6), Inches(1.8), Inches(5.8), Inches(4.9)
    )
    _place_images(slide, ctx, image_pane)
    _add_bullets(slide, ctx, text_pane, count=ctx.rng.randint(3, 5), size=20)


def _place_images(slide: object, ctx: SlideContext, pane: Box) -> None:
    """Fill ``pane`` with one or two stacked placeholder images."""
    count = ctx.rng.randint(1, 2)
    if count == 1:
        slide.shapes.add_picture(
            ctx.placeholder_image(720, 720),
            pane.left,
            pane.top,
            pane.width,
            pane.height,
        )
        return
    gap = Inches(0.2)
    each_h = (pane.height - gap) // 2
    for i in range(2):
        top = pane.top + i * (each_h + gap)
        slide.shapes.add_picture(
            ctx.placeholder_image(720, 480), pane.left, top, pane.width, each_h
        )


def chart_slide(slide: object, ctx: SlideContext, index: int) -> None:
    style.add_heading(slide, _HEADING, ctx.heading_text(index), ctx.theme)
    chart_data = CategoryChartData()
    chart_data.categories = [ctx.lorem.words(1) for _ in range(5)]
    for _ in range(ctx.rng.randint(1, 3)):
        chart_data.add_series(
            ctx.lorem.title(2),
            tuple(ctx.rng.randint(1, 100) for _ in range(5)),
        )
    box = Box(Inches(1.5), Inches(1.7), Inches(10.3), Inches(5.2))
    slide.shapes.add_chart(
        ctx.rng.choice(_CHART_TYPES),
        box.left,
        box.top,
        box.width,
        box.height,
        chart_data,
    )


def table_slide(slide: object, ctx: SlideContext, index: int) -> None:
    style.add_heading(slide, _HEADING, ctx.heading_text(index), ctx.theme)
    box = Box(Inches(0.8), Inches(1.7), Inches(11.7), Inches(5.0))
    cols = ctx.rng.randint(3, 5)
    row_h = Pt(30)
    # Fit rows to the box height so the table never overflows.
    max_rows = max(2, int(box.height / int(row_h)))
    rows = min(max_rows, ctx.rng.randint(4, 9))
    shape = slide.shapes.add_table(rows, cols, box.left, box.top, box.width, box.height)
    table = shape.table
    for c in range(cols):
        table.cell(0, c).text = ctx.lorem.words(1).title()
    for r in range(1, rows):
        for c in range(cols):
            table.cell(r, c).text = ctx.lorem.words(1)
    style.style_table(table, ctx.theme, font_size=12, row_height=row_h)


def video_slide(slide: object, ctx: SlideContext, index: int) -> None:
    style.add_heading(slide, _HEADING, ctx.heading_text(index), ctx.theme)
    poster_box = Box(Inches(2.4), Inches(1.8), Inches(8.5), Inches(4.4))
    poster = _video_poster(ctx)
    picture = slide.shapes.add_picture(
        poster, poster_box.left, poster_box.top, poster_box.width, poster_box.height
    )
    picture.click_action.hyperlink.address = ctx.video_url

    caption = slide.shapes.add_textbox(
        Inches(2.4), Inches(6.4), Inches(8.5), Inches(0.6)
    )
    run = caption.text_frame.paragraphs[0].add_run()
    run.text = "▶  Watch on YouTube"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.name = ctx.theme.body_font
    run.font.color.rgb = ctx.theme.primary
    run.hyperlink.address = ctx.video_url


# --- helpers ------------------------------------------------------------


def _video_poster(ctx: SlideContext) -> io.BytesIO:
    img = Image.new(
        "RGB",
        (960, 540),
        (
            ctx.theme.divider_background[0],
            ctx.theme.divider_background[1],
            ctx.theme.divider_background[2],
        ),
    )
    draw = ImageDraw.Draw(img)
    cx, cy, r = 480, 270, 70
    draw.ellipse([cx - r, cy - r, cx + r, cy + r], fill=(255, 255, 255))
    draw.polygon(
        [(cx - 22, cy - 36), (cx - 22, cy + 36), (cx + 40, cy)],
        fill=(ctx.theme.primary[0], ctx.theme.primary[1], ctx.theme.primary[2]),
    )
    return _png(img)


def _png(image: Image.Image) -> io.BytesIO:
    stream = io.BytesIO()
    image.save(stream, format="PNG")
    stream.seek(0)
    return stream
